import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9 Widescreen standard
prs.slide_height = Inches(7.5)

# Color Scheme Definition
COLOR_BG = RGBColor(17, 24, 39)       # Dark Slate Gray / Off-black #111827
COLOR_RED = RGBColor(239, 68, 68)     # Accent Crimson / Neon Red #EF4444
COLOR_CYAN = RGBColor(6, 182, 212)    # Accent Teal / Cyber Cyan #06B6D4
COLOR_WHITE = RGBColor(255, 255, 255) # Pure White
COLOR_SILVER = RGBColor(156, 163, 175)# Secondary Text Silver #9CA3AF
COLOR_CARD = RGBColor(31, 41, 55)     # Card Background #1F2937

# Helper function to set slide background
def set_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

# Helper function to create title section on content slides
def add_slide_header(slide, title_text, category="VULNERABILITY TECHNICAL BRIEFING"):
    # Add a thin cyan category label
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = "Trebuchet MS"
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_CYAN
    
    # Add the main title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Trebuchet MS"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE

# Helper to create bullet points in a modern look
def add_bullet_points(slide, left, top, width, height, points):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    for i, pt in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(12)
        
        # Check if point is a sub-bullet or header
        if pt.startswith("  - "):
            p.text = "   •  " + pt[4:]
            p.font.name = "Calibri"
            p.font.size = Pt(15)
            p.font.color.rgb = COLOR_SILVER
        elif pt.startswith("[CRITICAL]"):
            p.text = pt.replace("[CRITICAL]", "⚠️ ")
            p.font.name = "Calibri"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = COLOR_RED
        elif pt.startswith("[HIGHLIGHT]"):
            p.text = pt.replace("[HIGHLIGHT]", "⚡ ")
            p.font.name = "Calibri"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = COLOR_CYAN
        else:
            p.text = "•  " + pt
            p.font.name = "Calibri"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            
    return box

# Helper to draw a stylized code box
def add_code_box(slide, left, top, width, height, lines):
    # Add background card for code
    card = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD
    card.line.color.rgb = COLOR_CYAN
    card.line.width = Pt(1)
    
    # Add text inside the card
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(4)
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.LEFT
        
        # Color code annotations or highlights if they match keyword patterns
        if "PROT_READ" in line or "MAP_PRIVATE" in line:
            p.font.color.rgb = COLOR_CYAN
        elif "madvise" in line or "/proc/self/mem" in line:
            p.font.color.rgb = COLOR_RED
        elif line.startswith("//") or line.startswith("/*") or line.endswith("*/"):
            p.font.color.rgb = COLOR_SILVER
            
    return card

# Helper to draw visual info cards
def add_info_card(slide, left, top, width, height, title, subtitle, content, border_color):
    card = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    # Title
    p_title = tf.paragraphs[0]
    p_title.text = title.upper()
    p_title.font.name = "Trebuchet MS"
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    p_title.space_after = Pt(4)
    
    # Subtitle
    if subtitle:
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = "Calibri"
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = border_color
        p_sub.space_after = Pt(10)
        
    # Content
    for text_line in content:
        p_cont = tf.add_paragraph()
        p_cont.text = "• " + text_line
        p_cont.font.name = "Calibri"
        p_cont.font.size = Pt(13)
        p_cont.font.color.rgb = COLOR_SILVER
        p_cont.space_after = Pt(6)

# ==========================================
# SLIDE 1: Title Slide
# ==========================================
slide_layout = prs.slide_layouts[6] # Blank Layout
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)

# Title Text Box
title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True

# Main title
p1 = tf.paragraphs[0]
p1.text = "DIRTY COW EXPLOIT"
p1.font.name = "Trebuchet MS"
p1.font.size = Pt(60)
p1.font.bold = True
p1.font.color.rgb = COLOR_RED
p1.space_after = Pt(6)

p2 = tf.add_paragraph()
p2.text = "Deep-Dive CVE-2016-5195 & Container Security Lab"
p2.font.name = "Trebuchet MS"
p2.font.size = Pt(28)
p2.font.color.rgb = COLOR_WHITE
p2.space_after = Pt(18)

p3 = tf.add_paragraph()
p3.text = "A Complete End-to-End Demonstration and Privilege Escalation Walkthrough on Docker"
p3.font.name = "Calibri"
p3.font.size = Pt(14)
p3.font.color.rgb = COLOR_CYAN

# Bottom metadata
meta_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8))
tf_meta = meta_box.text_frame
p_meta = tf_meta.paragraphs[0]
p_meta.text = "TECHNICAL ASSESSMENT & LAB HANDS-ON DEMONSTRATION  |  SYSTEM SECURITY SERIES"
p_meta.font.name = "Trebuchet MS"
p_meta.font.size = Pt(11)
p_meta.font.bold = True
p_meta.font.color.rgb = COLOR_SILVER


# ==========================================
# SLIDE 2: Agenda
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Briefing & Laboratory Agenda", "EXECUTIVE PRESENTATION GUIDE")

# Add layout splits
add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), [
    "01. Executive Overview & Timeline",
    "  - Key definitions, discovery context, and CVSS severity metrics.",
    "02. Linux Kernel Memory Subsystems",
    "  - Virtual addressing, paging memory layouts, and the Copy-On-Write design.",
    "03. Mechanics of the Vulnerability",
    "  - Deep dive into the race condition, madvise(), and /proc/self/mem writes.",
    "04. Exploitation inside Containers",
    "  - Container-host kernel sharing architecture and threat models."
])

add_bullet_points(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), [
    "05. Hands-on Docker Lab Showcase",
    "  - Code walkthrough of dirtyc0w.c and Docker environment architecture.",
    "06. Execution & Diagnostic Walkthrough",
    "  - Step-by-step reproduction guide and interpreting kernel safety outcomes.",
    "07. Real-world Exploits & Impacts",
    "  - Explaining Android root, container breakout vectors, and server escapes.",
    "08. Defence, Mitigation & Remediation",
    "  - Kernel patch overview, container hardening configurations, and isolation rules."
])


# ==========================================
# SLIDE 3: Executive Summary (The Threat)
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Executive Summary: CVE-2016-5195", "VULNERABILITY TECHNICAL BRIEFING")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Core Vulnerability Definition",
    "  - A kernel-level race condition in Linux's Copy-on-Write (COW) memory subsystem.",
    "  - Bypasses traditional memory access control, allowing read-only files to be overwritten.",
    "The Severity Metric",
    "  - CVSS v3 Severity Score: 7.8 (High) / CVSS v2 Score: 10.0 (Critical).",
    "  - Highly reliable, zero complex requirements, zero-privilege prerequisites.",
    "[CRITICAL] Impact Matrix",
    "  - Grants local privilege escalation to root in seconds.",
    "  - Leads to full container breakout if host kernels are vulnerable.",
    "Wide-ranging Prevalence",
    "  - Existed in the kernel for 9 years (since v2.6.22) before discovery in Oct 2016.",
    "  - Impacted virtually every Linux distro, Android, and shared container server."
])

# Highlights right card box
add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "Threat Landscape Facts", "Vulnerability Timeline & Intel",
    [
        "Discovered by security researcher Phil Oester.",
        "Discovered actively exploited in-the-wild via packet capture analysis.",
        "Root cause resides in the memory manager's 'get_user_pages' implementation (mm/gup.c).",
        "Because virtual servers share a single host OS kernel, a single tenant could take control of the physical server."
    ],
    COLOR_RED
)


# ==========================================
# SLIDE 4: Linux Memory Management Foundations
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Linux Memory Management Foundations", "TECHNICAL DEEP DIVE")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Virtual Memory Architecture",
    "  - Applications run in clean virtual address spaces, isolated from the physical hardware.",
    "  - Memory pages (typically 4KB) map virtual addresses to physical RAM frames.",
    "Page Fault Mechanism",
    "  - When a process requests a page not currently residing in physical RAM, the CPU fires a page fault trap.",
    "  - The Kernel fault-handler steps in to load the page from disk/swap into RAM.",
    "Process Protection Flags",
    "  - Pages contain permissions: Read (R), Write (W), Execute (X).",
    "  - Read-only mappings (like shared system libraries) trigger a Segmentation Fault if writes are attempted."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "Memory Isolation Rules", "Core Kernel Rules",
    [
        "Process Separation: Process A cannot read or write Process B's memory space.",
        "Kernel Space Protection: User space programs can never directly modify kernel memory rings.",
        "Shared Memory Optimization: Multiple processes reading the exact same binaries share physical memory to save space."
    ],
    COLOR_CYAN
)


# ==========================================
# SLIDE 5: The Copy-On-Write (COW) Concept
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "The Copy-On-Write (COW) Design", "MEMORY OPTIMIZATION")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "The Efficiency Problem",
    "  - When a process forks, copying all memory pages instantly is expensive and slow.",
    "The Copy-On-Write Optimization",
    "  - Parent and child initially point to the identical physical page structures.",
    "  - The page flags are marked as read-only for both processes.",
    "Triggers on Write Action",
    "  - When either process attempts to write to the page, the CPU traps a page fault.",
    "  - The kernel intercepts, creates an independent copy of the page in physical memory, updates the writing process's page table, and marks it as writable.",
    "[HIGHLIGHT] Safe & Efficient",
    "  - Modifying process gets its private modified page; other processes continue safely sharing the original unmodified page cache."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "The Normal COW Flow", "Kernel Fault Sequencer",
    [
        "Step 1: Write request on read-only shared page.",
        "Step 2: CPU raises Page Fault exception.",
        "Step 3: Kernel allocates a new physical frame.",
        "Step 4: Kernel copies original data to the new frame.",
        "Step 5: Updates Page Table entries to the new frame.",
        "Step 6: Flushes TLB cache, changes permission to Writeable, and completes write."
    ],
    COLOR_CYAN
)


# ==========================================
# SLIDE 6: The Vulnerability (Race Condition)
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "The Vulnerability: COW Broken", "CRITICAL ANALYSIS")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Breaking the COW Sequence",
    "  - The kernel handles COW page table changes in multiple separate, non-atomic steps.",
    "  - If another system call changes page states in the middle of a COW process, the kernel's state machine loses track.",
    "The Attacker's Interruption",
    "  - An attacker starts a write process to a private mapping of a read-only file.",
    "  - Simultaneously, a second thread aggressively calls madvise(..., MADV_DONTNEED).",
    "[CRITICAL] Kernel Desynchronization",
    "  - madvise strips away the mapping just as the write handler finishes the page fault, causing the kernel to mistake the original, read-only cache page for the new COW page.",
    "  - Result: The write bypasses all permissions and modifies the immutable file on disk."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "Vulnerability Mechanics", "Get User Pages Race",
    [
        "During a write, the kernel calls get_user_pages() to find the page frame.",
        "The first search finds the read-only page and determines COW is required (returns flag FOLL_WRITE missing).",
        "The kernel copies the page, but madvise(MADV_DONTNEED) frees the new page tables.",
        "A retry occurs: the kernel fails to see the new page, fallback gets the original page, and incorrectly writes to it with write flags still active."
    ],
    COLOR_RED
)


# ==========================================
# SLIDE 7: Key System Calls
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Key System Calls & Components", "VULNERABILITY PREREQUISITES")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "mmap() System Call",
    "  - Maps files directly into a process's virtual memory region.",
    "  - Exploit uses MAP_PRIVATE: changes should be local and not write back to disk.",
    "madvise(..., MADV_DONTNEED)",
    "  - Tells the kernel: 'I do not need this memory range anymore; release associated resources and free the backing physical pages.'",
    "  - Forces the next read/write to look up the original page tables.",
    "/proc/self/mem Interface",
    "  - Virtual file representing the process's own virtual memory.",
    "  - Writing here allows a program to bypass some hardware-level CPU write-protection registers, shifting validation onto kernel software check logic."
])

add_code_box(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), [
    "// 1. Private mapping of read-only file",
    "map = mmap(NULL, size, PROT_READ,",
    "           MAP_PRIVATE, fd, 0);",
    "",
    "// 2. Thread A: Discard mapped pages",
    "void *madviseThread(void *arg) {",
    "  while(1) {",
    "    madvise(map, 100, MADV_DONTNEED);",
    "  }",
    "}",
    "",
    "// 3. Thread B: Write via procfs",
    "void *writeThread(void *arg) {",
    "  int f = open(\"/proc/self/mem\", O_RDWR);",
    "  while(1) {",
    "    lseek(f, map, SEEK_SET);",
    "    write(f, \"PAYLOAD\", 7);",
    "  }",
    "}"
])


# ==========================================
# SLIDE 8: Step-by-Step Race Condition Flow
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Step-by-Step Race Condition Flow", "THE EXPLOIT SEQUENCE")

# We will create a timeline layout using 3 vertical info cards side by side
card_width = Inches(3.6)
card_height = Inches(4.8)

add_info_card(
    slide, Inches(0.8), Inches(1.8), card_width, card_height,
    "1. Initiation Phase", "Threads Setup",
    [
        "Process maps 'target.txt' (Root-owned, Read-only) using mmap() privately.",
        "Thread A loops continuously calling madvise(map, 100, MADV_DONTNEED).",
        "Thread B loops calling write() on '/proc/self/mem' at the mapped address."
    ],
    COLOR_CYAN
)

add_info_card(
    slide, Inches(4.8), Inches(1.8), card_width, card_height,
    "2. The Intercept Race", "Kernel Processing",
    [
        "Thread B triggers a page fault because mapping is marked read-only.",
        "Kernel allocates a new physical frame and performs a Copy-On-Write (COW).",
        "Before Thread B updates the page table pointers, Thread A executes madvise().",
        "Kernel frees the page table entries for the new COW copy immediately."
    ],
    COLOR_RED
)

add_info_card(
    slide, Inches(8.8), Inches(1.8), card_width, card_height,
    "3. The Hijack", "Direct Disk Write",
    [
        "Kernel retries write loop, but the page tables are missing.",
        "Kernel assumes COW was already processed, resolving to the original cached page.",
        "Thread B writes the payload directly into the original page cache.",
        "Kernel flushes dirty cache to disk, permanently altering the read-only target."
    ],
    COLOR_RED
)


# ==========================================
# SLIDE 9: Dirty COW in Docker
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Dirty COW inside Docker Containers", "CONTAINER SECURITY INSIGHTS")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "The Kernel Sharing Reality",
    "  - Docker containers share the host Linux kernel directly.",
    "  - Containers are NOT hypervisor-isolated virtual machines.",
    "Isolation Boundaries Bypassed",
    "  - Traditional Docker namespaces separate filesystems, users, and networks.",
    "  - But a kernel flaw like Dirty COW breaks the core hypervisor-like assumption.",
    "[CRITICAL] Container Vulnerability Scope",
    "  - Even if the container has strict root-mapping isolation (User Namespaces disabled), an unprivileged user inside can exploit the kernel.",
    "  - The exploit directly corrupts physical page cache owned by the host operating system."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "Host vs. Guest Kernel Sharing", "Architectural Truths",
    [
        "Kernel Shared: Yes. Bypassing namespace separation completely.",
        "Host Filesystem Exposure: If host-level files are mapped privately, the container user can modify files outside their boundary.",
        "Privilege Escalation: Low-privilege container user -> Root-level container user -> Host system root control."
    ],
    COLOR_RED
)


# ==========================================
# SLIDE 10: Container Escape vs. Privilege Escalation
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Privilege Escalation to Container Escape", "THREAT PATHWAY ANALYSIS")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Attack Stage 1: Container Escalation",
    "  - Attacker gets low-privilege shell inside a container (e.g. web app exploit).",
    "  - Compiles and runs Dirty COW targeting `/etc/passwd` or `/etc/shadow`.",
    "  - Modifies root credentials inside the container to gain container root.",
    "Attack Stage 2: Container Breakout",
    "  - Container root searches for shared host volumes or mapped host files.",
    "  - Attacker overwrites host-owned binaries (e.g., cron jobs, systemd services) exposed via mount points.",
    "Attack Stage 3: Full Host Takeover",
    "  - The host executes the modified system-level binary as root.",
    "  - Reverse shell triggers, giving the attacker root command-line access to the host server."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "Common Attack Pathways", "Exploitation Flowchart",
    [
        "1. Web Exploit (RCE) -> Unprivileged Shell",
        "2. Dirty COW -> Write root-owned read-only configuration inside container.",
        "3. Local Privilege Escalation -> Root in Container.",
        "4. Escape: Modify shared Unix socket (e.g., docker.sock) or overwrite host binaries via mounts.",
        "5. Complete Host Infrastructure Compromise."
    ],
    COLOR_RED
)


# ==========================================
# SLIDE 11: Exploit Implementation (C Code)
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Exploit Implementation: C Code Breakdown", "CODE ANALYSIS")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Exploit Core Structure",
    "  - Uses the native POSIX threads library (`pthread.h`) for concurrent race loops.",
    "The Mapping Setup",
    "  - File mapped with `MAP_PRIVATE` and `PROT_READ` options.",
    "  - Guarantees the system will reject normal writes at compile/run time.",
    "The Thread Race Loop",
    "  - Runs millions of iterations of `madvise` and `write` to guarantee intercepting a page fault.",
    "[HIGHLIGHT] Simple yet Devastating",
    "  - Less than 120 lines of C code are needed to reliably compromise standard operating systems."
])

add_code_box(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), [
    "// Compiling the exploit with pthread support",
    "gcc -pthread dirtyc0w.c -o dirtyc0w",
    "",
    "// Core C Race Synchronization Block",
    "void *procSelfMemThread(void *arg) {",
    "  char *str = (char *)arg;",
    "  int f = open(\"/proc/self/mem\", O_RDWR);",
    "  for(int i = 0; i < 10000000; i++) {",
    "    lseek(f, (uintptr_t)map, SEEK_SET);",
    "    write(f, str, strlen(str));",
    "  }",
    "}",
    "",
    "// Target Execution",
    "./dirtyc0w /etc/passwd \"root:patched...\""
])


# ==========================================
# SLIDE 12: Hands-on Lab Environment
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Hands-on Lab: Container Environment", "DOCKER LAB INFRASTRUCTURE")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Isolated Lab Architecture",
    "  - Container base: Ubuntu 20.04 equipped with standard compile tools.",
    "  - Includes a low-privilege user named `victim` (UID: 1000).",
    "Immutable Target Configuration",
    "  - A file `/var/secret/target.txt` owned by `root:root` with permissions `444` (Read-only for all).",
    "  - Contains the secure flag: 'THIS_IS_A_SECRET_KEY_THAT_IS_READ_ONLY'.",
    "Demonstration Script (`run_demo.sh`)",
    "  - Automates environment diagnosis.",
    "  - Compiles and runs the exploit C code as the unprivileged user.",
    "  - Evaluates changes and outputs vulnerability health results."
])

add_code_box(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), [
    "FROM ubuntu:20.04",
    "RUN apt-get update && apt-get install -y gcc make",
    "",
    "# Create read-only, root-owned file",
    "RUN mkdir -p /var/secret && \\",
    "    echo \"SECRET_KEY\" > /var/secret/target.txt && \\",
    "    chmod 444 /var/secret/target.txt && \\",
    "    chown root:root /var/secret/target.txt",
    "",
    "# Setup restricted user",
    "RUN useradd -m -s /bin/bash victim",
    "COPY dirtyc0w.c /home/victim/dirtyc0w.c",
    "COPY run_demo.sh /home/victim/run_demo.sh",
    "USER victim",
    "WORKDIR /home/victim",
    "CMD [\"/bin/bash\", \"run_demo.sh\"]"
])


# ==========================================
# SLIDE 13: Executing the Exploit & Interpretation
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Executing the Lab & Outcome Analysis", "PRACTICAL LAB SHOWCASE")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Step 1: Building the Container Image",
    "  - Command: docker build -t dirtycow-lab .",
    "Step 2: Launching the Lab Container",
    "  - Command: docker run --rm -it dirtycow-lab",
    "Possible Outcome A: Attack Successful",
    "  - If the host OS runs a pre-patch Linux kernel (before Oct 2016).",
    "  - Result: Target file successfully overwritten to 'DIRTY_COW_SUCCESS!!!'.",
    "Possible Outcome B: Attack Prevented / Patched",
    "  - If the host runs a modern patched kernel (e.g., standard Windows WSL2 or modern Linux host).",
    "  - Result: Exploit runs but target remains safe; kernel correctly serializes COW memory writes."
])

# Code output preview block on the right
add_code_box(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), [
    "=== Dirty COW (CVE-2016-5195) Container Lab ===",
    "",
    "[*] System Diagnostics:",
    "  - Current User: victim (UID: 1000)",
    "  - Host Linux Kernel: 5.15.133.1-microsoft-standard-WSL2",
    "",
    "[*] Examining Target File Properties:",
    "  - File Owner: root | Permissions: -r--r--r--",
    "  - Current Contents: THIS_IS_A_SECRET_KEY_THAT_IS_READ_ONLY",
    "",
    "[*] Compiling Exploit C Code...",
    "[+] Compilation successful! Output binary: ./dirtyc0w",
    "",
    "[*] Triggering Exploit...",
    "procSelfMemThread completed 10000000 writes",
    "",
    "[*] Verifying Attack Status:",
    "[+] ATTACK PREVENTED / MITIGATED (Patched Kernel)"
])


# ==========================================
# SLIDE 14: Real-world Exploitation & Cases
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Real-world Exploitation & Notable Cases", "SECURITY IMPACT ANALYSIS")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Mass Android Rooting Tools",
    "  - Dirty COW was quickly weaponized to root billions of Android devices.",
    "  - Allowed attackers to bypass Android's SELinux rules and gain system root permissions.",
    "Shared Hosting Infrastructure Risks",
    "  - Cloud providers hosting un-isolated VPS environments faced high risk.",
    "  - A single low-tier user could corrupt shared system libraries on the hypervisor host.",
    "Malware Integration",
    "  - Integrated into Linux botnets and advanced persistent threats (APTs).",
    "  - Used as standard post-exploitation tools to secure root backdoors."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "Weaponization History", "Exploitation Vectors",
    [
        "Android Root: Bypassed early Android security scopes natively without bootloader unlock.",
        "Malware: Integrated into ELF binaries for silent kernel-level privilege elevation.",
        "Container escapes: Exploited in shared multi-tenant clusters before container security runtimes matured."
    ],
    COLOR_RED
)


# ==========================================
# SLIDE 15: Defense & Mitigation Strategies
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Defense & Hardening Strategies", "REMEDIATION AND DEFENCE")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Primary Mitigation: Host Patching",
    "  - Always patch the underlying host Linux kernel. Since containers share the host kernel, patching the host automatically secures all running containers.",
    "Docker Security Best Practices",
    "  - Run containers with non-root default users (`USER victim`).",
    "  - Implement read-only container filesystems (`--read-only`) to limit disk write vectors.",
    "Security Profiles",
    "  - Enable AppArmor or SELinux profiles to restrict container interactions.",
    "  - Configure custom `seccomp` filters to restrict dangerous system calls if not required (e.g. blocking /proc modifications)."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "Container Hardening Rules", "Host Defense Guidelines",
    [
        "Kernel Patching: Update Linux hosts immediately.",
        "Minimal Runtimes: Use container profiles to restrict unnecessary privileges.",
        "Explicit Seccomp: Limit dangerous syscall access by default.",
        "Read-Only Containers: Mount filesystems as read-only to stop post-exploitation persistence."
    ],
    COLOR_CYAN
)


# ==========================================
# SLIDE 16: Conclusion & Key Takeaways
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Conclusion & Core Lessons", "TECHNICAL DEBRIEFING")

add_bullet_points(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Key Takeaway 1: Kernel Vulnerabilities are Shared",
    "  - Container isolation is only as secure as the shared host kernel.",
    "  - Kernel security bugs completely bypass container namespace boundaries.",
    "Key Takeaway 2: Race Conditions are Highly Lethal",
    "  - Non-atomic multi-step kernel operations create critical race conditions.",
    "  - Even simple logic can bypass strict physical hardware access controls.",
    "[HIGHLIGHT] Security Paradigm Shift",
    "  - Dirty COW forced a major evolution in container security runtimes, emphasizing host patching, seccomp filters, and hardened sandboxes (gVisor/Kata Containers)."
])

add_info_card(
    slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8),
    "The Golden Rule of Containers", "Final Recommendation",
    [
        "Containers are NOT isolation virtual machines.",
        "Always treat host kernel patch levels as the absolute foundation of container cluster security.",
        "Implement a defense-in-depth architecture to limit exposure if a kernel compromise occurs."
    ],
    COLOR_CYAN
)

# Save the presentation
output_path = "Dirty_COW_Presentation.pptx"
prs.save(output_path)
print(f"[+] Successfully generated presentation at: {os.path.abspath(output_path)}")
